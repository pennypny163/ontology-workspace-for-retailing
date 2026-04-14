"""
零售解决方案本体工作台 - 后端API服务
基于 FastAPI + SQLAlchemy + MySQL
"""
import os
import io
import json
import re
import tempfile
import traceback
import httpx
from typing import Optional, List, Dict, Any
from datetime import datetime

from fastapi import FastAPI, Query, HTTPException, Depends, UploadFile, File, Form
from fastapi.staticfiles import StaticFiles
from fastapi.responses import RedirectResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

from sqlalchemy import (
    create_engine, Column, Integer, String, Text, Enum, DECIMAL,
    JSON, DateTime, Boolean, ForeignKey, func, Index, text as db_text,
    case as sql_case
)
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, Session, relationship

# ============================================================
# 数据库配置（通过环境变量获取）
# ============================================================
DATABASE_URL = os.environ.get("DATABASE_URL")
if not DATABASE_URL:
    MYSQL_HOST = os.environ.get("MYSQL_HOST")
    MYSQL_PORT = os.environ.get("MYSQL_PORT")
    MYSQL_DATABASE = os.environ.get("MYSQL_DATABASE")
    MYSQL_USER = os.environ.get("MYSQL_USER")
    MYSQL_PASSWORD = os.environ.get("MYSQL_PASSWORD")

    if MYSQL_HOST or MYSQL_DATABASE or MYSQL_USER or MYSQL_PASSWORD:
        DB_HOST = MYSQL_HOST or "127.0.0.1"
        DB_PORT = MYSQL_PORT or "3306"
        DB_NAME = MYSQL_DATABASE or "ontology"
        DB_USER = MYSQL_USER or "root"
        DB_PASS = MYSQL_PASSWORD or ""
        DATABASE_URL = f"mysql+pymysql://{DB_USER}:{DB_PASS}@{DB_HOST}:{DB_PORT}/{DB_NAME}?charset=utf8mb4"
    else:
        DATABASE_URL = "sqlite:///./ontology.db"

# 示例本地配置：
# 如果你希望使用 SQLite 进行本地开发，可直接启动后端而不设置数据库环境变量，应用会自动使用 ./ontology.db。
# 如果你希望使用 MySQL，可设置下面变量：
# export MYSQL_HOST=127.0.0.1
# export MYSQL_PORT=3306
# export MYSQL_DATABASE=ontology
# export MYSQL_USER=root
# export MYSQL_PASSWORD=""
# 或者直接：
# export DATABASE_URL="mysql+pymysql://root:password@127.0.0.1:3306/ontology?charset=utf8mb4"

# ============================================================
# 大模型 API 配置
# ============================================================
LLM_API_KEY = os.environ.get("LLM_API_KEY", "")
LLM_API_BASE = os.environ.get("LLM_API_BASE", "")
# 使用默认模型或通过环境变量覆盖
LLM_MODEL = os.environ.get("LLM_MODEL", "")

engine = create_engine(DATABASE_URL, pool_pre_ping=True, pool_recycle=3600, echo=False)
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

# ============================================================
# SQLAlchemy 模型定义
# ============================================================

class AssetType(Base):
    """资产类型表"""
    __tablename__ = "asset_types"

    id = Column(Integer, primary_key=True, autoincrement=True)
    code = Column(String(50), unique=True, nullable=False, comment="英文唯一编码")
    name = Column(String(100), nullable=False, comment="中文名称")
    description = Column(Text, comment="类型说明")
    icon_name = Column(String(50), comment="前端图标标识")
    color = Column(String(20), comment="主题色")
    display_order = Column(Integer, default=0, comment="排序")
    is_active = Column(Boolean, default=True, comment="是否启用")
    created_at = Column(DateTime, default=datetime.now)
    updated_at = Column(DateTime, default=datetime.now, onupdate=datetime.now)


class OntologyObject(Base):
    """资产对象表（本体节点）"""
    __tablename__ = "ontology_objects"

    id = Column(Integer, primary_key=True, autoincrement=True)
    object_type_code = Column(String(50), ForeignKey("asset_types.code"), nullable=False, comment="关联asset_types.code")
    name = Column(String(200), nullable=False, comment="资产名称")
    short_description = Column(String(500), comment="列表页短描述")
    full_description = Column(Text, comment="详情页完整描述")
    status = Column(String(20), default="draft", comment="发布状态: draft/reviewed/published")
    ai_confidence = Column(DECIMAL(5, 2), comment="AI置信度 0~100")
    tags_json = Column(JSON, comment="标签数组")
    attachment_filename = Column(String(300), comment="附件原始文件名")
    attachment_path = Column(String(500), comment="附件存储路径")
    attachment_size = Column(Integer, comment="附件大小(字节)")
    external_link = Column(String(1000), comment="外部链接（如腾讯云能力介绍页URL）")
    kpi_formula = Column(String(500), comment="KPI计算公式（可选）")
    created_by = Column(String(100), comment="创建者")
    created_at = Column(DateTime, default=datetime.now)
    updated_at = Column(DateTime, default=datetime.now, onupdate=datetime.now)
    is_deleted = Column(Boolean, default=False, comment="软删除")

    # 关系
    asset_type = relationship("AssetType", foreign_keys=[object_type_code], primaryjoin="OntologyObject.object_type_code == AssetType.code")


class OntologyRelation(Base):
    """对象关系表（本体边）"""
    __tablename__ = "ontology_relations"

    id = Column(Integer, primary_key=True, autoincrement=True)
    from_object_id = Column(Integer, ForeignKey("ontology_objects.id", ondelete="CASCADE"), nullable=False)
    relation_type = Column(String(50), nullable=False, comment="关系类型")
    to_object_id = Column(Integer, ForeignKey("ontology_objects.id", ondelete="CASCADE"), nullable=False)
    remark = Column(String(500), comment="备注")
    display_order = Column(Integer, default=0, comment="排序")
    created_at = Column(DateTime, default=datetime.now)
    updated_at = Column(DateTime, default=datetime.now, onupdate=datetime.now)
    status = Column(String(20), default="active", comment="关系状态 active/inactive/pending")
    ai_confidence = Column(DECIMAL(5, 2), comment="AI置信度 0~100")
    weight = Column(DECIMAL(5, 2), default=1.0, comment="关系权重 0~10")
    created_by = Column(String(100), comment="创建者")
    is_deleted = Column(Boolean, default=False, comment="软删除")

    from_object = relationship("OntologyObject", foreign_keys=[from_object_id])
    to_object = relationship("OntologyObject", foreign_keys=[to_object_id])


class ObjectProperty(Base):
    """对象扩展属性表"""
    __tablename__ = "object_properties"

    id = Column(Integer, primary_key=True, autoincrement=True)
    object_id = Column(Integer, ForeignKey("ontology_objects.id", ondelete="CASCADE"), nullable=False)
    property_key = Column(String(100), nullable=False, comment="属性键")
    property_value = Column(Text, comment="属性值")
    created_at = Column(DateTime, default=datetime.now)
    updated_at = Column(DateTime, default=datetime.now, onupdate=datetime.now)


class RelationType(Base):
    """关系类型配置表（语义层）"""
    __tablename__ = "relation_types"

    id = Column(Integer, primary_key=True, autoincrement=True)
    code = Column(String(50), unique=True, nullable=False, comment="关系类型编码")
    name_zh = Column(String(100), nullable=False, comment="中文名称")
    name_en = Column(String(100), comment="英文名称")
    reverse_name_zh = Column(String(100), comment="反向中文名称")
    description = Column(Text, comment="描述")
    is_directed = Column(Boolean, default=True, comment="是否有向")
    color = Column(String(20), default="#6366f1", comment="连线颜色")
    line_style = Column(String(20), default="solid", comment="线型")
    display_order = Column(Integer, default=0, comment="排序")
    is_active = Column(Boolean, default=True, comment="是否启用")
    created_at = Column(DateTime, default=datetime.now)
    updated_at = Column(DateTime, default=datetime.now, onupdate=datetime.now)


class RelationSchemaRule(Base):
    """关系模式约束规则表（语义层）"""
    __tablename__ = "relation_schema_rules"

    id = Column(Integer, primary_key=True, autoincrement=True)
    source_type_code = Column(String(50), nullable=False, comment="源类型code")
    relation_type_code = Column(String(50), nullable=False, comment="关系类型code")
    target_type_code = Column(String(50), nullable=False, comment="目标类型code")
    is_allowed = Column(Boolean, default=True, comment="是否允许")
    validation_level = Column(String(20), default="warning", comment="校验级别")
    remark = Column(String(500), comment="备注")
    created_at = Column(DateTime, default=datetime.now)
    updated_at = Column(DateTime, default=datetime.now, onupdate=datetime.now)


class AssetTypePropertySchema(Base):
    """资产类型属性Schema定义表（语义层）"""
    __tablename__ = "asset_type_property_schema"

    id = Column(Integer, primary_key=True, autoincrement=True)
    object_type_code = Column(String(50), nullable=False, comment="资产类型code")
    property_key = Column(String(100), nullable=False, comment="属性键")
    property_name = Column(String(200), comment="属性中文名")
    property_type = Column(String(50), default="text", comment="属性值类型")
    is_required = Column(Boolean, default=False, comment="是否必填")
    is_filterable = Column(Boolean, default=False, comment="是否可筛选")
    is_displayable = Column(Boolean, default=True, comment="是否展示")
    default_value = Column(Text, comment="默认值")
    display_order = Column(Integer, default=0, comment="排序")
    created_at = Column(DateTime, default=datetime.now)
    updated_at = Column(DateTime, default=datetime.now, onupdate=datetime.now)


# ============================================================
# Pydantic 响应模型
# ============================================================

class AssetTypeSummary(BaseModel):
    id: int
    code: str
    name: str
    description: Optional[str] = None
    icon_name: Optional[str] = None
    color: Optional[str] = None
    display_order: int = 0
    record_count: int = 0
    sample_items: List[str] = []

class ObjectListItem(BaseModel):
    id: int
    object_type_code: str
    name: str
    short_description: Optional[str] = None
    status: str = "draft"
    ai_confidence: Optional[float] = None
    tags: List[str] = []
    created_by: Optional[str] = None
    created_at: Optional[str] = None
    updated_at: Optional[str] = None

class RelationItem(BaseModel):
    relation_type: str
    relation_label_cn: str
    related_object_id: int
    related_object_type: str
    related_object_type_name: str
    related_object_name: str

class ObjectDetail(BaseModel):
    id: int
    object_type_code: str
    object_type_name: str
    name: str
    short_description: Optional[str] = None
    full_description: Optional[str] = None
    status: str = "draft"
    ai_confidence: Optional[float] = None
    tags: List[str] = []
    attachment_filename: Optional[str] = None
    attachment_path: Optional[str] = None
    attachment_size: Optional[int] = None
    external_link: Optional[str] = None
    kpi_formula: Optional[str] = None
    created_by: Optional[str] = None
    created_at: Optional[str] = None
    updated_at: Optional[str] = None
    relations: List[RelationItem] = []

class AssetSummaryResponse(BaseModel):
    total_types: int
    total_objects: int
    published_count: int
    reviewed_count: int
    draft_count: int

class PaginatedObjects(BaseModel):
    items: List[ObjectListItem]
    total: int
    page: int
    page_size: int

# ============================================================
# 关系类型中文映射
# ============================================================
RELATION_TYPE_CN = {
    "serves_for": "服务于",
    "solves": "解决",
    "depends_on": "依赖",
    "measured_by": "由…衡量",
    "related_to": "相关",
    "belongs_to": "属于",
    "validated_by": "被…验证",
}

# ============================================================
# FastAPI 应用
# ============================================================
app = FastAPI(title="零售解决方案本体工作台 API", version="1.0.0")

# 上传附件存储目录
UPLOAD_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)


@app.on_event("startup")
def on_startup():
    """应用启动时自动补齐数据库字段"""
    with engine.connect() as conn:
        # 检查并添加附件字段
        try:
            conn.execute(
                db_text("ALTER TABLE ontology_objects ADD COLUMN attachment_filename VARCHAR(300) COMMENT '附件原始文件名'")
            )
        except Exception:
            pass  # 字段已存在
        try:
            conn.execute(
                db_text("ALTER TABLE ontology_objects ADD COLUMN attachment_path VARCHAR(500) COMMENT '附件存储路径'")
            )
        except Exception:
            pass
        try:
            conn.execute(
                db_text("ALTER TABLE ontology_objects ADD COLUMN attachment_size INT COMMENT '附件大小(字节)'")
            )
        except Exception:
            pass
        try:
            conn.execute(
                db_text("ALTER TABLE ontology_objects ADD COLUMN external_link VARCHAR(1000) COMMENT '外部链接（如产品能力介绍页URL）'")
            )
        except Exception:
            pass
        try:
            conn.execute(
                db_text("ALTER TABLE ontology_objects ADD COLUMN kpi_formula VARCHAR(500) COMMENT 'KPI计算公式（可选）'")
            )
        except Exception:
            pass
        conn.commit()

# CORS 中间件
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def get_db():
    """获取数据库会话"""
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

# ============================================================
# 根路径重定向
# ============================================================
@app.get("/")
async def root():
    return RedirectResponse(url="/static/index.html")

# ============================================================
# API 路由
# ============================================================

@app.get("/api/asset-types", response_model=List[AssetTypeSummary])
def get_asset_types(db: Session = Depends(get_db)):
    """
    获取所有资产类型卡片数据（总览页）
    每个类型包含记录数和示例资产名称
    """
    types = db.query(AssetType).filter(AssetType.is_active == True).order_by(AssetType.display_order).all()
    
    result = []
    for t in types:
        # 统计该类型下的资产记录数（排除软删除）
        count = db.query(func.count(OntologyObject.id)).filter(
            OntologyObject.object_type_code == t.code,
            OntologyObject.is_deleted == False
        ).scalar()
        
        # 获取前3个示例资产名称
        samples = db.query(OntologyObject.name).filter(
            OntologyObject.object_type_code == t.code,
            OntologyObject.is_deleted == False
        ).order_by(OntologyObject.id).limit(3).all()
        sample_names = [s[0] for s in samples]
        
        result.append(AssetTypeSummary(
            id=t.id,
            code=t.code,
            name=t.name,
            description=t.description,
            icon_name=t.icon_name,
            color=t.color,
            display_order=t.display_order,
            record_count=count,
            sample_items=sample_names,
        ))
    
    return result


@app.get("/api/asset-types/{type_code}/objects", response_model=PaginatedObjects)
def get_objects_by_type(
    type_code: str,
    keyword: str = Query("", description="按名称搜索"),
    page: int = Query(1, ge=1, description="页码"),
    page_size: int = Query(20, ge=1, le=100, description="每页条数"),
    db: Session = Depends(get_db)
):
    """
    获取某个类型下的资产列表，支持分页和关键词搜索
    """
    # 验证类型是否存在
    asset_type = db.query(AssetType).filter(AssetType.code == type_code).first()
    if not asset_type:
        raise HTTPException(status_code=404, detail=f"资产类型 {type_code} 不存在")
    
    query = db.query(OntologyObject).filter(
        OntologyObject.object_type_code == type_code,
        OntologyObject.is_deleted == False
    )
    
    if keyword:
        query = query.filter(OntologyObject.name.like(f"%{keyword}%"))
    
    total = query.count()
    objects = query.order_by(OntologyObject.id).offset((page - 1) * page_size).limit(page_size).all()
    
    items = []
    for obj in objects:
        tags = obj.tags_json if obj.tags_json else []
        items.append(ObjectListItem(
            id=obj.id,
            object_type_code=obj.object_type_code,
            name=obj.name,
            short_description=obj.short_description,
            status=obj.status,
            ai_confidence=float(obj.ai_confidence) if obj.ai_confidence else None,
            tags=tags,
            created_by=obj.created_by,
            created_at=obj.created_at.strftime("%Y-%m-%d %H:%M") if obj.created_at else None,
            updated_at=obj.updated_at.strftime("%Y-%m-%d %H:%M") if obj.updated_at else None,
        ))
    
    return PaginatedObjects(items=items, total=total, page=page, page_size=page_size)


@app.get("/api/objects/{object_id}", response_model=ObjectDetail)
def get_object_detail(object_id: int, db: Session = Depends(get_db)):
    """
    获取单个资产详情（右侧抽屉）
    包含基础信息、标签、关联资产列表
    """
    obj = db.query(OntologyObject).filter(
        OntologyObject.id == object_id,
        OntologyObject.is_deleted == False
    ).first()
    
    if not obj:
        raise HTTPException(status_code=404, detail="资产不存在")
    
    # 获取类型中文名
    asset_type = db.query(AssetType).filter(AssetType.code == obj.object_type_code).first()
    type_name = asset_type.name if asset_type else obj.object_type_code
    
    # 获取该对象的所有关系（包括正向和反向）
    relations = []
    
    # 正向关系：from_object_id = object_id
    forward_rels = db.query(OntologyRelation).filter(
        OntologyRelation.from_object_id == object_id
    ).order_by(OntologyRelation.display_order).all()
    
    for rel in forward_rels:
        to_obj = db.query(OntologyObject).filter(OntologyObject.id == rel.to_object_id).first()
        if to_obj and not to_obj.is_deleted:
            to_type = db.query(AssetType).filter(AssetType.code == to_obj.object_type_code).first()
            relations.append(RelationItem(
                relation_type=rel.relation_type,
                relation_label_cn=RELATION_TYPE_CN.get(rel.relation_type, rel.relation_type),
                related_object_id=to_obj.id,
                related_object_type=to_obj.object_type_code,
                related_object_type_name=to_type.name if to_type else to_obj.object_type_code,
                related_object_name=to_obj.name,
            ))
    
    # 反向关系：to_object_id = object_id
    reverse_rels = db.query(OntologyRelation).filter(
        OntologyRelation.to_object_id == object_id
    ).order_by(OntologyRelation.display_order).all()
    
    # 反向关系文案映射
    reverse_label_map = {
        "serves_for": "被…服务",
        "solves": "被…解决",
        "depends_on": "被…依赖",
        "measured_by": "衡量…",
        "related_to": "相关",
        "belongs_to": "包含",
        "validated_by": "验证…",
    }
    
    for rel in reverse_rels:
        from_obj = db.query(OntologyObject).filter(OntologyObject.id == rel.from_object_id).first()
        if from_obj and not from_obj.is_deleted:
            from_type = db.query(AssetType).filter(AssetType.code == from_obj.object_type_code).first()
            relations.append(RelationItem(
                relation_type=rel.relation_type,
                relation_label_cn=reverse_label_map.get(rel.relation_type, rel.relation_type),
                related_object_id=from_obj.id,
                related_object_type=from_obj.object_type_code,
                related_object_type_name=from_type.name if from_type else from_obj.object_type_code,
                related_object_name=from_obj.name,
            ))
    
    tags = obj.tags_json if obj.tags_json else []
    
    return ObjectDetail(
        id=obj.id,
        object_type_code=obj.object_type_code,
        object_type_name=type_name,
        name=obj.name,
        short_description=obj.short_description,
        full_description=obj.full_description,
        status=obj.status,
        ai_confidence=float(obj.ai_confidence) if obj.ai_confidence else None,
        tags=tags,
        attachment_filename=obj.attachment_filename,
        attachment_path=obj.attachment_path,
        attachment_size=obj.attachment_size,
        external_link=obj.external_link,
        kpi_formula=obj.kpi_formula,
        created_by=obj.created_by,
        created_at=obj.created_at.strftime("%Y-%m-%d %H:%M") if obj.created_at else None,
        updated_at=obj.updated_at.strftime("%Y-%m-%d %H:%M") if obj.updated_at else None,
        relations=relations,
    )


@app.get("/api/asset-summary", response_model=AssetSummaryResponse)
def get_asset_summary(db: Session = Depends(get_db)):
    """
    获取资产总览统计信息
    """
    total_types = db.query(func.count(AssetType.id)).filter(AssetType.is_active == True).scalar()
    
    base_query = db.query(OntologyObject).filter(OntologyObject.is_deleted == False)
    total_objects = base_query.count()
    published_count = base_query.filter(OntologyObject.status == "published").count()
    reviewed_count = base_query.filter(OntologyObject.status == "reviewed").count()
    draft_count = base_query.filter(OntologyObject.status == "draft").count()
    
    return AssetSummaryResponse(
        total_types=total_types,
        total_objects=total_objects,
        published_count=published_count,
        reviewed_count=reviewed_count,
        draft_count=draft_count,
    )


# ============================================================
# 文件上传 + 大模型提取 API
# ============================================================

def parse_file_content(file_bytes: bytes, filename: str) -> str:
    """
    根据文件类型解析文件内容为文本
    支持 PPTX、DOCX、PDF、TXT
    """
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    text_parts = []

    if ext in ("pptx",):
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation(io.BytesIO(file_bytes))
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                # 提取普通文本框/标题/正文
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                # 提取表格内容（PPT中表格是重要信息载体）
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_texts.append(cell_text)
                        if row_texts:
                            slide_texts.append(" | ".join(row_texts))
                # 提取分组形状中的文本（SmartArt等会被解析为grouped shapes）
                if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                    try:
                        for child_shape in shape.shapes:
                            if hasattr(child_shape, "text") and child_shape.text.strip():
                                slide_texts.append(child_shape.text.strip())
                    except Exception:
                        pass
            if slide_texts:
                text_parts.append(f"[幻灯片 {slide_idx}]\n" + "\n".join(slide_texts))

    elif ext in ("docx",):
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())

    elif ext in ("pdf",):
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page_idx, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_parts.append(f"[页面 {page_idx}]\n{page_text.strip()}")

    elif ext in ("txt", "md"):
        text_parts.append(file_bytes.decode("utf-8", errors="ignore"))

    else:
        raise ValueError(f"不支持的文件格式: .{ext}，请上传 PPTX、DOCX、PDF 或 TXT 文件")

    full_text = "\n\n".join(text_parts)
    # 截断过长文本（8000字符可覆盖大多数PPT的全部幻灯片文本）
    max_chars = 8000
    if len(full_text) > max_chars:
        full_text = full_text[:max_chars] + f"\n\n[...文本过长已截断，共 {len(full_text)} 字符，已保留前 {max_chars} 字符...]"

    return full_text


async def call_llm_extract(text_content: str, db: Session) -> Dict[str, Any]:
    """
    调用大模型 API 从文本中提取资产对象
    使用流式(stream)接收，避免长时间无响应被网关断开
    """
    # 获取当前所有资产类型code列表（精简格式）
    types = db.query(AssetType).filter(AssetType.is_active == True).all()
    type_codes = [f"{t.code}({t.name})" for t in types]
    types_str = ",".join(type_codes)

    # 精简版 prompt —— 减少 token 消耗和生成时间
    system_prompt = f"""从文档提取零售方案资产,只输出JSON。
可用type: {types_str}
关系type: serves_for,solves,depends_on,measured_by,related_to,belongs_to,validated_by
格式:{{"objects":[{{"object_type_code":"","name":"","short_description":"","tags":[],"ai_confidence":85}}],"relations":[{{"from_name":"","to_name":"","relation_type":"","remark":""}}],"summary":""}}
尽量完整提取文档中的所有资产对象和关系,最多10个对象10个关系,只输JSON,无需解释。"""

    user_prompt = text_content

    headers = {
        "Authorization": f"Bearer {LLM_API_KEY}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": LLM_MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        "temperature": 0.2,
        "max_tokens": 2500,
        "stream": True,  # 使用流式接收，防止网关超时
    }

    # 使用流式请求：网关只在"无数据传输"时才会触发超时
    # 流式模式下每生成一个 token 都会有数据返回，不会触发60秒网关超时
    content_parts = []
    timeout_cfg = httpx.Timeout(connect=10.0, read=120.0, write=10.0, pool=10.0)
    
    async with httpx.AsyncClient(timeout=timeout_cfg) as client:
        async with client.stream(
            "POST",
            f"{LLM_API_BASE}/chat/completions",
            headers=headers,
            json=payload,
        ) as resp:
            if resp.status_code != 200:
                body = await resp.aread()
                raise HTTPException(
                    status_code=502,
                    detail=f"大模型 API 调用失败: HTTP {resp.status_code} - {body.decode('utf-8', errors='ignore')[:500]}"
                )
            # 逐行读取 SSE 流
            async for line in resp.aiter_lines():
                line = line.strip()
                if not line or not line.startswith("data:"):
                    continue
                data_str = line[len("data:"):].strip()
                if data_str == "[DONE]":
                    break
                try:
                    chunk = json.loads(data_str)
                    delta = chunk.get("choices", [{}])[0].get("delta", {})
                    if "content" in delta and delta["content"]:
                        content_parts.append(delta["content"])
                except (json.JSONDecodeError, IndexError, KeyError):
                    continue

    content = "".join(content_parts)
    
    if not content.strip():
        raise HTTPException(status_code=500, detail="大模型返回内容为空")

    # 尝试提取 JSON（大模型可能返回 markdown 代码块包裹的 JSON）
    json_match = re.search(r"```(?:json)?\s*\n?([\s\S]*?)\n?```", content)
    if json_match:
        json_str = json_match.group(1)
    else:
        json_str = content.strip()

    # 处理 qwen3.5 可能在 JSON 前输出思考过程（<think>...</think>）
    think_match = re.search(r"</think>\s*([\s\S]*)", json_str)
    if think_match:
        json_str = think_match.group(1).strip()

    try:
        extracted = json.loads(json_str)
    except json.JSONDecodeError:
        raise HTTPException(
            status_code=500,
            detail=f"大模型返回的内容无法解析为JSON: {content[:500]}"
        )

    return extracted


# --- Pydantic 模型：上传提取 ---
class ExtractedObject(BaseModel):
    object_type_code: str
    name: str
    short_description: Optional[str] = None
    full_description: Optional[str] = None
    tags: List[str] = []
    ai_confidence: Optional[float] = None
    external_link: Optional[str] = None
    kpi_formula: Optional[str] = None

class ExtractedRelation(BaseModel):
    from_name: str
    to_name: str
    relation_type: str
    remark: Optional[str] = None

class ExtractionResult(BaseModel):
    objects: List[ExtractedObject] = []
    relations: List[ExtractedRelation] = []
    summary: Optional[str] = None
    source_filename: str = ""
    text_preview: str = ""

class ConfirmSaveRequest(BaseModel):
    objects: List[ExtractedObject]
    relations: List[ExtractedRelation]
    created_by: Optional[str] = None


@app.post("/api/upload-extract")
async def upload_and_extract(
    file: UploadFile = File(..., description="上传文件（PPT/DOC/PDF/TXT）"),
    db: Session = Depends(get_db),
):
    """
    上传文件并通过大模型提取资产对象
    使用 SSE (Server-Sent Events) 流式返回，防止网关超时
    依次发送: progress 事件 → result/error 事件
    """
    # 检查文件类型
    allowed_extensions = {"pptx", "docx", "pdf", "txt", "md"}
    filename = file.filename or "unknown"
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext not in allowed_extensions:
        raise HTTPException(status_code=400, detail=f"不支持的文件格式 .{ext}，请上传 PPTX、DOCX、PDF 或 TXT")

    # 检查文件大小（最大 50MB）
    file_bytes = await file.read()
    if len(file_bytes) > 50 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="文件大小超过 50MB 限制")

    import asyncio

    async def event_stream():
        """SSE 事件流生成器"""
        # 进度: 解析文件
        yield f"data: {json.dumps({'type': 'progress', 'step': 'parsing', 'message': '正在解析文件内容...'}, ensure_ascii=False)}\n\n"
        await asyncio.sleep(0.05)  # 确保数据被刷出

        try:
            text_content = parse_file_content(file_bytes, filename)
        except ValueError as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"
            return
        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': f'文件解析失败: {str(e)}'}, ensure_ascii=False)}\n\n"
            return

        if not text_content.strip():
            yield f"data: {json.dumps({'type': 'error', 'message': '文件内容为空，无法提取'}, ensure_ascii=False)}\n\n"
            return

        # 进度: 调用大模型
        yield f"data: {json.dumps({'type': 'progress', 'step': 'extracting', 'message': f'已提取 {len(text_content)} 字文本，正在调用 AI 模型提取资产...'}, ensure_ascii=False)}\n\n"
        await asyncio.sleep(0.05)

        # 每5秒发送心跳，防止网关超时
        extraction_done = False
        extracted = None
        extract_error = None

        async def do_extract():
            nonlocal extracted, extract_error, extraction_done
            try:
                extracted = await call_llm_extract(text_content, db)
            except HTTPException as e:
                extract_error = e.detail
            except httpx.TimeoutException:
                extract_error = "大模型响应超时，请尝试上传更小的文件。"
            except Exception as e:
                traceback.print_exc()
                extract_error = f"大模型提取失败: {str(e)}"
            finally:
                extraction_done = True

        # 启动提取任务
        extract_task = asyncio.create_task(do_extract())

        # 同时发送心跳，每4秒发一次，保持连接活跃
        heartbeat_count = 0
        while not extraction_done:
            await asyncio.sleep(4)
            if not extraction_done:
                heartbeat_count += 1
                yield f"data: {json.dumps({'type': 'progress', 'step': 'extracting', 'message': f'AI 正在分析中... ({heartbeat_count * 4}秒)'}, ensure_ascii=False)}\n\n"

        await extract_task  # 确保任务完成

        if extract_error:
            yield f"data: {json.dumps({'type': 'error', 'message': extract_error}, ensure_ascii=False)}\n\n"
            return

        # 进度: 组装结果
        yield f"data: {json.dumps({'type': 'progress', 'step': 'done', 'message': '提取完成，正在组装结果...'}, ensure_ascii=False)}\n\n"
        await asyncio.sleep(0.05)

        # 组装结果
        objects = []
        for obj in extracted.get("objects", []):
            objects.append({
                "object_type_code": obj.get("object_type_code", ""),
                "name": obj.get("name", ""),
                "short_description": obj.get("short_description"),
                "full_description": obj.get("full_description"),
                "tags": obj.get("tags", []),
                "ai_confidence": obj.get("ai_confidence"),
            })

        relations = []
        for rel in extracted.get("relations", []):
            relations.append({
                "from_name": rel.get("from_name", ""),
                "to_name": rel.get("to_name", ""),
                "relation_type": rel.get("relation_type", "related_to"),
                "remark": rel.get("remark"),
            })

        text_preview = text_content[:500] + ("..." if len(text_content) > 500 else "")

        result = {
            "type": "result",
            "data": {
                "objects": objects,
                "relations": relations,
                "summary": extracted.get("summary", ""),
                "source_filename": filename,
                "text_preview": text_preview,
            }
        }
        yield f"data: {json.dumps(result, ensure_ascii=False)}\n\n"

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",  # 告诉 Nginx 不要缓冲
        }
    )


@app.post("/api/confirm-save")
def confirm_save_assets(
    req: ConfirmSaveRequest,
    db: Session = Depends(get_db),
):
    """
    用户确认提取结果后，将资产对象和关系写入数据库
    所有资产初始状态为 draft（草稿）
    """
    # 验证所有 object_type_code 是否有效
    valid_codes = {t.code for t in db.query(AssetType.code).filter(AssetType.is_active == True).all()}
    for obj in req.objects:
        if obj.object_type_code not in valid_codes:
            raise HTTPException(
                status_code=400,
                detail=f"无效的资产类型: {obj.object_type_code}"
            )

    # 创建资产对象
    name_to_id: Dict[str, int] = {}
    created_objects = []
    for obj in req.objects:
        new_obj = OntologyObject(
            object_type_code=obj.object_type_code,
            name=obj.name,
            short_description=obj.short_description,
            full_description=obj.full_description,
            status="draft",
            ai_confidence=obj.ai_confidence,
            tags_json=obj.tags,
            external_link=obj.external_link,
            kpi_formula=obj.kpi_formula,
            created_by=req.created_by or "AI提取",
            created_at=datetime.now(),
            updated_at=datetime.now(),
        )
        db.add(new_obj)
        db.flush()  # 获取自增ID
        name_to_id[obj.name] = new_obj.id
        created_objects.append({
            "id": new_obj.id,
            "name": obj.name,
            "type": obj.object_type_code,
        })

    # 创建关系
    created_relations = []
    for rel in req.relations:
        from_id = name_to_id.get(rel.from_name)
        to_id = name_to_id.get(rel.to_name)
        if from_id and to_id:
            new_rel = OntologyRelation(
                from_object_id=from_id,
                relation_type=rel.relation_type,
                to_object_id=to_id,
                remark=rel.remark,
                created_at=datetime.now(),
            )
            db.add(new_rel)
            created_relations.append({
                "from": rel.from_name,
                "to": rel.to_name,
                "type": rel.relation_type,
            })

    db.commit()

    return {
        "success": True,
        "message": f"成功保存 {len(created_objects)} 个资产对象和 {len(created_relations)} 个关系",
        "created_objects": created_objects,
        "created_relations": created_relations,
    }


# ============================================================
# 引导表单 API：获取已有资产列表（用于建立关系选择）
# ============================================================

@app.get("/api/existing-objects")
def get_existing_objects(
    keyword: str = Query("", description="搜索关键词"),
    type_code: str = Query("", description="按类型过滤"),
    limit: int = Query(50, ge=1, le=200),
    db: Session = Depends(get_db),
):
    """获取已有资产对象列表，用于引导表单中建立关系时选择"""
    query = db.query(OntologyObject).filter(OntologyObject.is_deleted == False)
    if type_code:
        query = query.filter(OntologyObject.object_type_code == type_code)
    if keyword:
        query = query.filter(OntologyObject.name.like(f"%{keyword}%"))
    objects = query.order_by(OntologyObject.id.desc()).limit(limit).all()

    result = []
    for obj in objects:
        at = db.query(AssetType).filter(AssetType.code == obj.object_type_code).first()
        result.append({
            "id": obj.id,
            "name": obj.name,
            "object_type_code": obj.object_type_code,
            "object_type_name": at.name if at else obj.object_type_code,
            "status": obj.status,
        })
    return result


# ============================================================
# 引导表单 API：AI 增强（补充描述和标签）
# ============================================================

class AIEnhanceRequest(BaseModel):
    object_type_code: str
    name: str
    short_description: Optional[str] = None
    tags: List[str] = []


@app.post("/api/ai-enhance")
async def ai_enhance_form(req: AIEnhanceRequest, db: Session = Depends(get_db)):
    """
    用大模型对用户填写的表单进行增强：
    - 自动补充完整描述 full_description
    - 自动推荐标签 tags
    - 自动给出置信度 ai_confidence
    使用流式接收避免超时
    """
    at = db.query(AssetType).filter(AssetType.code == req.object_type_code).first()
    type_name = at.name if at else req.object_type_code

    system_prompt = f"""你是零售方案资产助手。用户通过表单新建了一个"{type_name}"类型的资产，请帮助增强信息。
只输出JSON，格式如下：
{{"full_description":"详细描述(100-200字)","suggested_tags":["标签1","标签2","标签3","标签4","标签5"],"ai_confidence":85}}
- full_description: 根据名称和简述，生成专业且详尽的描述，涵盖该资产在零售行业中的应用场景、核心价值
- suggested_tags: 推荐5个相关标签，简短精炼(2-6字)，涵盖行业、场景、技术等维度
- ai_confidence: 根据信息完整度给出置信度(0-100)
只输出JSON，不要解释。"""

    user_content = f"资产类型: {type_name}\n名称: {req.name}"
    if req.short_description:
        user_content += f"\n简述: {req.short_description}"
    if req.tags:
        user_content += f"\n已有标签: {', '.join(req.tags)}"

    headers = {
        "Authorization": f"Bearer {LLM_API_KEY}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": LLM_MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_content},
        ],
        "temperature": 0.3,
        "max_tokens": 800,
        "stream": True,
    }

    content_parts = []
    timeout_cfg = httpx.Timeout(connect=10.0, read=120.0, write=10.0, pool=10.0)

    try:
        async with httpx.AsyncClient(timeout=timeout_cfg) as client:
            async with client.stream(
                "POST",
                f"{LLM_API_BASE}/chat/completions",
                headers=headers,
                json=payload,
            ) as resp:
                if resp.status_code != 200:
                    body = await resp.aread()
                    raise HTTPException(
                        status_code=502,
                        detail=f"AI增强失败: HTTP {resp.status_code} - {body.decode('utf-8', errors='ignore')[:300]}"
                    )
                async for line in resp.aiter_lines():
                    line = line.strip()
                    if not line or not line.startswith("data:"):
                        continue
                    data_str = line[len("data:"):].strip()
                    if data_str == "[DONE]":
                        break
                    try:
                        chunk = json.loads(data_str)
                        delta = chunk.get("choices", [{}])[0].get("delta", {})
                        if "content" in delta and delta["content"]:
                            content_parts.append(delta["content"])
                    except (json.JSONDecodeError, IndexError, KeyError):
                        continue
    except httpx.TimeoutException:
        raise HTTPException(status_code=504, detail="AI增强超时，请稍后重试")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"AI增强失败: {str(e)}")

    content = "".join(content_parts)
    if not content.strip():
        raise HTTPException(status_code=500, detail="AI返回内容为空")

    # 提取JSON
    json_match = re.search(r"```(?:json)?\s*\n?([\s\S]*?)\n?```", content)
    json_str = json_match.group(1) if json_match else content.strip()
    # 处理 <think> 标签
    think_match = re.search(r"</think>\s*([\s\S]*)", json_str)
    if think_match:
        json_str = think_match.group(1).strip()

    try:
        enhanced = json.loads(json_str)
    except json.JSONDecodeError:
        raise HTTPException(status_code=500, detail=f"AI返回内容无法解析: {content[:300]}")

    return {
        "full_description": enhanced.get("full_description", ""),
        "suggested_tags": enhanced.get("suggested_tags", []),
        "ai_confidence": enhanced.get("ai_confidence", 70),
    }


# ============================================================
# 引导表单 API：附件上传（独立上传附件，返回存储信息）
# ============================================================

@app.post("/api/form-upload-attachment")
async def form_upload_attachment(
    file: UploadFile = File(..., description="资产附件文件"),
):
    """
    独立的附件上传接口：
    - 用户在引导表单中为案例资产/数据资产/演示资产等上传附件文件
    - 文件存储到服务端 uploads/ 目录
    - 返回文件的存储路径、文件名、大小，供后续 form-submit 关联
    """
    import uuid as uuid_mod
    filename = file.filename or "unknown"

    # 检查文件大小（最大 100MB）
    file_bytes = await file.read()
    file_size = len(file_bytes)
    if file_size > 100 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="文件大小超过 100MB 限制")

    # 生成唯一存储路径，保留原始扩展名
    ext = filename.rsplit(".", 1)[-1] if "." in filename else ""
    stored_name = f"{uuid_mod.uuid4().hex}.{ext}" if ext else uuid_mod.uuid4().hex
    stored_path = os.path.join(UPLOAD_DIR, stored_name)

    with open(stored_path, "wb") as f:
        f.write(file_bytes)

    return {
        "success": True,
        "filename": filename,
        "stored_path": stored_name,  # 只存储相对文件名
        "file_size": file_size,
    }


@app.get("/api/attachments/{stored_name}")
async def download_attachment(stored_name: str):
    """下载/预览附件文件"""
    from fastapi.responses import FileResponse
    file_path = os.path.join(UPLOAD_DIR, stored_name)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="附件不存在")
    return FileResponse(file_path, filename=stored_name)


# ============================================================
# 引导表单 API：表单提交（创建资产 + 可选关系 + 可选附件）
# ============================================================

class FormRelation(BaseModel):
    target_object_id: int
    relation_type: str
    remark: Optional[str] = None


class FormSubmitRequest(BaseModel):
    object_type_code: str
    name: str
    short_description: Optional[str] = None
    full_description: Optional[str] = None
    tags: List[str] = []
    ai_confidence: Optional[float] = None
    relations: List[FormRelation] = []
    created_by: Optional[str] = None
    attachment_filename: Optional[str] = None
    attachment_path: Optional[str] = None
    attachment_size: Optional[int] = None
    external_link: Optional[str] = None
    kpi_formula: Optional[str] = None


@app.post("/api/form-submit")
def form_submit_asset(req: FormSubmitRequest, db: Session = Depends(get_db)):
    """
    引导表单提交接口：创建单个资产对象 + 与已有资产的关系 + 关联附件
    """
    # 验证类型
    at = db.query(AssetType).filter(AssetType.code == req.object_type_code, AssetType.is_active == True).first()
    if not at:
        raise HTTPException(status_code=400, detail=f"无效的资产类型: {req.object_type_code}")

    # 对于需要附件的类型，检查是否上传了附件
    FILE_REQUIRED_TYPES = {"case_asset", "demo_asset", "data_asset", "playbook"}
    if req.object_type_code in FILE_REQUIRED_TYPES and not req.attachment_path:
        raise HTTPException(
            status_code=400,
            detail=f"「{at.name}」类型的资产需要上传附件文件，请先上传相关文件"
        )

    # 创建对象
    new_obj = OntologyObject(
        object_type_code=req.object_type_code,
        name=req.name,
        short_description=req.short_description,
        full_description=req.full_description,
        status="draft",
        ai_confidence=req.ai_confidence,
        tags_json=req.tags if req.tags else [],
        attachment_filename=req.attachment_filename,
        attachment_path=req.attachment_path,
        attachment_size=req.attachment_size,
        external_link=req.external_link,
        kpi_formula=req.kpi_formula,
        created_by=req.created_by or "表单贡献",
        created_at=datetime.now(),
        updated_at=datetime.now(),
    )
    db.add(new_obj)
    db.flush()

    # 创建关系（当前对象 → 已有对象）
    created_relations = []
    for rel in req.relations:
        target_obj = db.query(OntologyObject).filter(
            OntologyObject.id == rel.target_object_id,
            OntologyObject.is_deleted == False
        ).first()
        if target_obj:
            new_rel = OntologyRelation(
                from_object_id=new_obj.id,
                relation_type=rel.relation_type,
                to_object_id=rel.target_object_id,
                remark=rel.remark,
                created_at=datetime.now(),
            )
            db.add(new_rel)
            created_relations.append({
                "from": req.name,
                "to": target_obj.name,
                "type": rel.relation_type,
            })

    db.commit()

    return {
        "success": True,
        "message": f"资产「{req.name}」已保存为草稿" + (f"（含附件 {req.attachment_filename}）" if req.attachment_filename else ""),
        "created_object": {
            "id": new_obj.id,
            "name": req.name,
            "type": req.object_type_code,
            "type_name": at.name,
        },
        "created_relations": created_relations,
    }


# ============================================================
# 自然语言多轮对话 API：与大模型交互，最多3轮追问
# ============================================================

class NLChatMessage(BaseModel):
    role: str  # "user" | "assistant"
    content: str

class NLChatRequest(BaseModel):
    messages: List[NLChatMessage]
    round_index: int = 1  # 当前是第几轮 (1~3)

class NLChatResponse(BaseModel):
    reply: str  # AI 回复文本
    is_final: bool = False  # 是否为最终轮（输出结构化结果）
    extraction: Optional[Dict[str, Any]] = None  # 最终轮时附带结构化提取结果


@app.post("/api/nl-chat", response_model=NLChatResponse)
async def nl_chat(req: NLChatRequest, db: Session = Depends(get_db)):
    """
    自然语言多轮对话接口：
    - 第1~2轮：AI根据用户描述进行追问，帮助补全资产信息
    - 第3轮：AI输出最终结构化资产 JSON
    使用 SSE 流式读取大模型响应防止超时
    """
    MAX_ROUNDS = 3
    round_index = req.round_index

    # 获取当前所有资产类型
    types = db.query(AssetType).filter(AssetType.is_active == True).all()
    type_codes = [f"{t.code}({t.name})" for t in types]
    types_str = ", ".join(type_codes)

    is_final = round_index >= MAX_ROUNDS

    if is_final:
        # 最终轮：要求输出结构化 JSON
        system_prompt = f"""你是零售解决方案资产助手。用户通过自然语言描述了一个零售方案，你已经与用户进行了多轮对话来收集信息。
现在请根据对话历史，提取并输出结构化的资产对象。

可用资产类型: {types_str}
可用关系类型: serves_for(服务于), solves(解决), depends_on(依赖), measured_by(由…衡量), related_to(相关), belongs_to(属于), validated_by(被…验证)

请严格按以下JSON格式输出，不要输出任何解释文字：
```json
{{
  "summary": "对用户方案的一句话总结",
  "objects": [
    {{
      "object_type_code": "类型code",
      "name": "资产名称",
      "short_description": "简短描述",
      "full_description": "详细描述(80-150字)",
      "tags": ["标签1", "标签2", "标签3"],
      "ai_confidence": 85
    }}
  ],
  "relations": [
    {{
      "from_name": "源资产名",
      "to_name": "目标资产名",
      "relation_type": "关系类型code",
      "remark": "备注"
    }}
  ]
}}
```
最多提取6个对象和6个关系。只输出JSON代码块，不要任何其他文字。"""
    else:
        # 追问轮：根据用户信息进行追问
        remaining = MAX_ROUNDS - round_index
        system_prompt = f"""你是零售解决方案资产助手，正在帮助用户通过自然语言描述来创建结构化资产。

可用资产类型: {types_str}

你的任务：
1. 分析用户已经提供的信息
2. 找出还缺少哪些关键信息（如目标行业/场景、服务的角色/用户群体、要解决的痛点、涉及的能力模块、相关KPI指标等）
3. 用友好、专业的语气提出2-3个有针对性的追问

注意事项：
- 追问要简洁明了，每个问题独占一行，用数字编号
- 不要重复用户已经提供的信息
- 追问内容要围绕零售行业方案的关键维度
- 当前还有 {remaining} 次追问机会（含本次），之后系统将生成最终结构化资产
- 先简要总结你从对话中已经了解到的信息，再提出追问
- 不要输出JSON，只输出自然语言"""

    # 构建消息列表
    llm_messages = [{"role": "system", "content": system_prompt}]
    for msg in req.messages:
        llm_messages.append({"role": msg.role, "content": msg.content})

    # 调用大模型（流式接收）
    headers = {
        "Authorization": f"Bearer {LLM_API_KEY}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": LLM_MODEL,
        "messages": llm_messages,
        "temperature": 0.4 if not is_final else 0.2,
        "max_tokens": 2000 if is_final else 800,
        "stream": True,
    }

    content_parts = []
    timeout_cfg = httpx.Timeout(connect=10.0, read=120.0, write=10.0, pool=10.0)

    try:
        async with httpx.AsyncClient(timeout=timeout_cfg) as client:
            async with client.stream(
                "POST",
                f"{LLM_API_BASE}/chat/completions",
                headers=headers,
                json=payload,
            ) as resp:
                if resp.status_code != 200:
                    body = await resp.aread()
                    raise HTTPException(
                        status_code=502,
                        detail=f"AI对话失败: HTTP {resp.status_code} - {body.decode('utf-8', errors='ignore')[:300]}"
                    )
                async for line in resp.aiter_lines():
                    line = line.strip()
                    if not line or not line.startswith("data:"):
                        continue
                    data_str = line[len("data:"):].strip()
                    if data_str == "[DONE]":
                        break
                    try:
                        chunk = json.loads(data_str)
                        delta = chunk.get("choices", [{}])[0].get("delta", {})
                        if "content" in delta and delta["content"]:
                            content_parts.append(delta["content"])
                    except (json.JSONDecodeError, IndexError, KeyError):
                        continue
    except httpx.TimeoutException:
        raise HTTPException(status_code=504, detail="AI对话超时，请稍后重试")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"AI对话失败: {str(e)}")

    content = "".join(content_parts)
    if not content.strip():
        raise HTTPException(status_code=500, detail="AI返回内容为空")

    # 处理 <think> 标签
    think_match = re.search(r"</think>\s*([\s\S]*)", content)
    if think_match:
        content = think_match.group(1).strip()

    # 如果是最终轮，尝试解析JSON
    extraction = None
    if is_final:
        json_match = re.search(r"```(?:json)?\s*\n?([\s\S]*?)\n?```", content)
        json_str = json_match.group(1) if json_match else content.strip()
        try:
            extraction = json.loads(json_str)
        except json.JSONDecodeError:
            # 如果解析失败，返回原文，让前端提示
            extraction = None

    return NLChatResponse(
        reply=content,
        is_final=is_final,
        extraction=extraction,
    )


# ============================================================
# 语义层 API
# ============================================================

@app.get("/api/semantic/relation-types")
def get_relation_types(db: Session = Depends(get_db)):
    """获取所有可用关系类型（从数据库配置读取，不再硬编码）"""
    types = db.query(RelationType).filter(RelationType.is_active == True).order_by(RelationType.display_order).all()
    return [{
        "id": t.id,
        "code": t.code,
        "name_zh": t.name_zh,
        "name_en": t.name_en,
        "reverse_name_zh": t.reverse_name_zh,
        "description": t.description,
        "is_directed": t.is_directed,
        "color": t.color,
        "line_style": t.line_style,
        "display_order": t.display_order,
    } for t in types]


@app.get("/api/semantic/schema-rules")
def get_schema_rules(db: Session = Depends(get_db)):
    """获取所有关系模式约束规则"""
    rules = db.query(RelationSchemaRule).filter(RelationSchemaRule.is_allowed == True).all()
    return [{
        "id": r.id,
        "source_type_code": r.source_type_code,
        "relation_type_code": r.relation_type_code,
        "target_type_code": r.target_type_code,
        "validation_level": r.validation_level,
        "remark": r.remark,
    } for r in rules]


@app.get("/api/semantic/object-type/{type_code}/property-schema")
def get_property_schema(type_code: str, db: Session = Depends(get_db)):
    """获取某资产类型的属性Schema定义"""
    schemas = db.query(AssetTypePropertySchema).filter(
        AssetTypePropertySchema.object_type_code == type_code
    ).order_by(AssetTypePropertySchema.display_order).all()
    return [{
        "id": s.id,
        "property_key": s.property_key,
        "property_name": s.property_name,
        "property_type": s.property_type,
        "is_required": s.is_required,
        "is_filterable": s.is_filterable,
        "is_displayable": s.is_displayable,
        "default_value": s.default_value,
        "display_order": s.display_order,
    } for s in schemas]


class ValidateRelationRequest(BaseModel):
    source_type_code: str
    relation_type_code: str
    target_type_code: str


@app.post("/api/semantic/validate-relation")
def validate_relation(req: ValidateRelationRequest, db: Session = Depends(get_db)):
    """校验关系三元组是否符合Schema规则"""
    rule = db.query(RelationSchemaRule).filter(
        RelationSchemaRule.source_type_code == req.source_type_code,
        RelationSchemaRule.relation_type_code == req.relation_type_code,
        RelationSchemaRule.target_type_code == req.target_type_code,
    ).first()

    if rule and rule.is_allowed:
        return {"valid": True, "validation_level": rule.validation_level, "message": "关系合法"}
    elif rule and not rule.is_allowed:
        return {"valid": False, "validation_level": rule.validation_level, "message": f"该关系被明确禁止: {rule.remark or ''}"}
    else:
        # 没有规则定义 -> 视为 warning
        return {"valid": False, "validation_level": "warning", "message": "该关系组合未在Schema中定义，请确认是否合理"}


# ============================================================
# Graph Projection API
# ============================================================

def _build_graph_node(obj: OntologyObject, asset_type_map: dict) -> dict:
    """构建前端 graph 节点数据"""
    at = asset_type_map.get(obj.object_type_code, {})
    return {
        "id": f"obj-{obj.id}",
        "objectId": obj.id,
        "label": obj.name,
        "type": obj.object_type_code,
        "typeName": at.get("name", obj.object_type_code),
        "color": at.get("color", "#6366f1"),
        "icon": at.get("icon_name", "Database"),
        "status": obj.status,
        "tags": obj.tags_json or [],
        "ai_confidence": float(obj.ai_confidence) if obj.ai_confidence else None,
        "short_description": obj.short_description,
    }


def _build_graph_edge(rel: OntologyRelation, rel_type_map: dict, schema_set: set) -> dict:
    """构建前端 graph 边数据"""
    rt = rel_type_map.get(rel.relation_type, {})
    # 从 ORM 关系中获取 source/target 类型
    from_type = rel.from_object.object_type_code if rel.from_object else ""
    to_type = rel.to_object.object_type_code if rel.to_object else ""
    schema_key = (from_type, rel.relation_type, to_type)
    return {
        "id": f"rel-{rel.id}",
        "relationId": rel.id,
        "source": f"obj-{rel.from_object_id}",
        "target": f"obj-{rel.to_object_id}",
        "relation_type": rel.relation_type,
        "label": rt.get("name_zh", rel.relation_type),
        "reverse_label": rt.get("reverse_name_zh", ""),
        "color": rt.get("color", "#6b7280"),
        "line_style": rt.get("line_style", "solid"),
        "is_schema_valid": schema_key in schema_set,
        "remark": rel.remark,
        "weight": float(rel.weight) if rel.weight else 1.0,
    }


def _get_asset_type_map(db: Session) -> dict:
    """获取资产类型映射表 {code: {name, color, icon_name}}"""
    types = db.query(AssetType).filter(AssetType.is_active == True).all()
    return {t.code: {"name": t.name, "color": t.color or "#6366f1", "icon_name": t.icon_name or "Database"} for t in types}


def _get_relation_type_map(db: Session) -> dict:
    """获取关系类型映射表 {code: {name_zh, reverse_name_zh, color, line_style}}"""
    types = db.query(RelationType).filter(RelationType.is_active == True).all()
    return {t.code: {
        "name_zh": t.name_zh, "reverse_name_zh": t.reverse_name_zh,
        "color": t.color or "#6b7280", "line_style": t.line_style or "solid",
    } for t in types}


def _get_schema_set(db: Session) -> set:
    """获取合法三元组集合 {(source_type, relation_type, target_type)}"""
    rules = db.query(RelationSchemaRule).filter(RelationSchemaRule.is_allowed == True).all()
    return {(r.source_type_code, r.relation_type_code, r.target_type_code) for r in rules}


@app.get("/api/graph/subgraph")
def get_subgraph(
    object_id: int = Query(..., description="中心节点ID"),
    depth: int = Query(1, ge=1, le=3, description="展开深度(1~3跳)"),
    db: Session = Depends(get_db),
):
    """
    获取某个节点的一跳/多跳子图
    返回前端可直接渲染的 nodes + edges 结构
    """
    asset_type_map = _get_asset_type_map(db)
    rel_type_map = _get_relation_type_map(db)
    schema_set = _get_schema_set(db)

    visited_ids = set()
    frontier = {object_id}
    all_object_ids = set()
    all_relation_ids = set()
    collected_rels = []

    for d in range(depth):
        if not frontier:
            break
        new_frontier = set()
        for oid in frontier:
            if oid in visited_ids:
                continue
            visited_ids.add(oid)
            all_object_ids.add(oid)

            # 正向关系
            fwd = db.query(OntologyRelation).filter(
                OntologyRelation.from_object_id == oid,
                OntologyRelation.is_deleted == False,
            ).all()
            for r in fwd:
                if r.id not in all_relation_ids:
                    # 检查 to_object 是否未删除
                    to_obj = db.query(OntologyObject).filter(OntologyObject.id == r.to_object_id, OntologyObject.is_deleted == False).first()
                    if to_obj:
                        all_relation_ids.add(r.id)
                        collected_rels.append(r)
                        all_object_ids.add(r.to_object_id)
                        new_frontier.add(r.to_object_id)

            # 反向关系
            rev = db.query(OntologyRelation).filter(
                OntologyRelation.to_object_id == oid,
                OntologyRelation.is_deleted == False,
            ).all()
            for r in rev:
                if r.id not in all_relation_ids:
                    from_obj = db.query(OntologyObject).filter(OntologyObject.id == r.from_object_id, OntologyObject.is_deleted == False).first()
                    if from_obj:
                        all_relation_ids.add(r.id)
                        collected_rels.append(r)
                        all_object_ids.add(r.from_object_id)
                        new_frontier.add(r.from_object_id)

        frontier = new_frontier - visited_ids

    # 构建节点
    objects = db.query(OntologyObject).filter(
        OntologyObject.id.in_(all_object_ids),
        OntologyObject.is_deleted == False,
    ).all()
    nodes = [_build_graph_node(obj, asset_type_map) for obj in objects]

    # 构建边
    edges = [_build_graph_edge(r, rel_type_map, schema_set) for r in collected_rels]

    return {
        "nodes": nodes,
        "edges": edges,
        "meta": {
            "center_object_id": object_id,
            "depth": depth,
            "total_nodes": len(nodes),
            "total_edges": len(edges),
        }
    }


@app.get("/api/graph/filter")
def get_filtered_graph(
    object_type_code: Optional[str] = Query(None, description="按资产类型筛选"),
    status: Optional[str] = Query(None, description="按状态筛选"),
    relation_type: Optional[str] = Query(None, description="按关系类型筛选"),
    tag: Optional[str] = Query(None, description="按标签筛选"),
    keyword: Optional[str] = Query(None, description="按关键词搜索"),
    limit: int = Query(50, ge=1, le=200, description="最大节点数"),
    db: Session = Depends(get_db),
):
    """
    按条件筛选图谱子集
    返回符合条件的 nodes + 它们之间的 edges
    """
    asset_type_map = _get_asset_type_map(db)
    rel_type_map = _get_relation_type_map(db)
    schema_set = _get_schema_set(db)

    # 筛选节点
    q = db.query(OntologyObject).filter(OntologyObject.is_deleted == False)
    if object_type_code:
        q = q.filter(OntologyObject.object_type_code == object_type_code)
    if status:
        q = q.filter(OntologyObject.status == status)
    if keyword:
        q = q.filter(OntologyObject.name.like(f"%{keyword}%"))
    if tag:
        q = q.filter(OntologyObject.tags_json.like(f'%"{tag}"%'))

    objects = q.order_by(OntologyObject.id.desc()).limit(limit).all()
    obj_ids = {obj.id for obj in objects}

    nodes = [_build_graph_node(obj, asset_type_map) for obj in objects]

    # 获取这些节点之间的边
    edge_q = db.query(OntologyRelation).filter(
        OntologyRelation.from_object_id.in_(obj_ids),
        OntologyRelation.to_object_id.in_(obj_ids),
        OntologyRelation.is_deleted == False,
    )
    if relation_type:
        edge_q = edge_q.filter(OntologyRelation.relation_type == relation_type)

    rels = edge_q.all()
    edges = [_build_graph_edge(r, rel_type_map, schema_set) for r in rels]

    return {
        "nodes": nodes,
        "edges": edges,
        "meta": {
            "total_nodes": len(nodes),
            "total_edges": len(edges),
            "filters_applied": {
                "object_type_code": object_type_code,
                "status": status,
                "relation_type": relation_type,
                "tag": tag,
                "keyword": keyword,
            }
        }
    }


@app.get("/api/graph/overview")
def get_graph_overview(
    limit: int = Query(80, ge=1, le=300, description="最大节点数"),
    db: Session = Depends(get_db),
):
    """
    获取图谱总览（默认首页展示）
    优先展示 published 和 reviewed 状态的资产
    """
    asset_type_map = _get_asset_type_map(db)
    rel_type_map = _get_relation_type_map(db)
    schema_set = _get_schema_set(db)

    # 按状态优先级取资产
    objects = db.query(OntologyObject).filter(
        OntologyObject.is_deleted == False,
    ).order_by(
        sql_case(
            (OntologyObject.status == "published", 0),
            (OntologyObject.status == "reviewed", 1),
            else_=2
        ),
        OntologyObject.id.desc()
    ).limit(limit).all()

    obj_ids = {obj.id for obj in objects}
    nodes = [_build_graph_node(obj, asset_type_map) for obj in objects]

    # 获取这些节点之间的所有关系
    rels = db.query(OntologyRelation).filter(
        OntologyRelation.from_object_id.in_(obj_ids),
        OntologyRelation.to_object_id.in_(obj_ids),
        OntologyRelation.is_deleted == False,
    ).all()
    edges = [_build_graph_edge(r, rel_type_map, schema_set) for r in rels]

    return {
        "nodes": nodes,
        "edges": edges,
        "meta": {
            "total_nodes": len(nodes),
            "total_edges": len(edges),
        }
    }


@app.get("/api/graph/relation-detail/{relation_id}")
def get_relation_detail(relation_id: int, db: Session = Depends(get_db)):
    """获取单个关系的详细信息（点击边时展示）"""
    rel = db.query(OntologyRelation).filter(
        OntologyRelation.id == relation_id,
        OntologyRelation.is_deleted == False,
    ).first()
    if not rel:
        raise HTTPException(status_code=404, detail="关系不存在")

    # 获取关系类型元信息
    rt = db.query(RelationType).filter(RelationType.code == rel.relation_type).first()

    # 获取两端对象
    from_obj = db.query(OntologyObject).filter(OntologyObject.id == rel.from_object_id).first()
    to_obj = db.query(OntologyObject).filter(OntologyObject.id == rel.to_object_id).first()

    # Schema 校验
    schema_valid = False
    if from_obj and to_obj:
        rule = db.query(RelationSchemaRule).filter(
            RelationSchemaRule.source_type_code == from_obj.object_type_code,
            RelationSchemaRule.relation_type_code == rel.relation_type,
            RelationSchemaRule.target_type_code == to_obj.object_type_code,
            RelationSchemaRule.is_allowed == True,
        ).first()
        schema_valid = rule is not None

    from_type = db.query(AssetType).filter(AssetType.code == from_obj.object_type_code).first() if from_obj else None
    to_type = db.query(AssetType).filter(AssetType.code == to_obj.object_type_code).first() if to_obj else None

    return {
        "id": rel.id,
        "relation_type": rel.relation_type,
        "name_zh": rt.name_zh if rt else rel.relation_type,
        "reverse_name_zh": rt.reverse_name_zh if rt else "",
        "color": rt.color if rt else "#6b7280",
        "line_style": rt.line_style if rt else "solid",
        "remark": rel.remark,
        "weight": float(rel.weight) if rel.weight else 1.0,
        "status": rel.status,
        "ai_confidence": float(rel.ai_confidence) if rel.ai_confidence else None,
        "is_schema_valid": schema_valid,
        "from_object": {
            "id": from_obj.id if from_obj else None,
            "name": from_obj.name if from_obj else "未知",
            "type_code": from_obj.object_type_code if from_obj else "",
            "type_name": from_type.name if from_type else "",
        },
        "to_object": {
            "id": to_obj.id if to_obj else None,
            "name": to_obj.name if to_obj else "未知",
            "type_code": to_obj.object_type_code if to_obj else "",
            "type_name": to_type.name if to_type else "",
        },
        "created_at": rel.created_at.strftime("%Y-%m-%d %H:%M") if rel.created_at else None,
        "created_by": rel.created_by,
    }


@app.get("/api/ontology/objects/{object_id}")
def get_ontology_object_detail(object_id: int, db: Session = Depends(get_db)):
    """
    增强版对象详情（用于 graph 右侧面板）
    除了基础信息，还包含：扩展属性、正向关系、反向关系、schema 校验
    """
    obj = db.query(OntologyObject).filter(
        OntologyObject.id == object_id,
        OntologyObject.is_deleted == False,
    ).first()
    if not obj:
        raise HTTPException(status_code=404, detail="资产不存在")

    at = db.query(AssetType).filter(AssetType.code == obj.object_type_code).first()
    rel_type_map = _get_relation_type_map(db)
    schema_set = _get_schema_set(db)

    # 扩展属性
    props = db.query(ObjectProperty).filter(ObjectProperty.object_id == object_id).all()
    properties = [{"key": p.property_key, "value": p.property_value} for p in props]

    # 正向关系
    fwd_rels = db.query(OntologyRelation).filter(
        OntologyRelation.from_object_id == object_id,
        OntologyRelation.is_deleted == False,
    ).all()
    forward_relations = []
    for r in fwd_rels:
        to_obj = db.query(OntologyObject).filter(OntologyObject.id == r.to_object_id, OntologyObject.is_deleted == False).first()
        if not to_obj:
            continue
        to_at = db.query(AssetType).filter(AssetType.code == to_obj.object_type_code).first()
        rt = rel_type_map.get(r.relation_type, {})
        sk = (obj.object_type_code, r.relation_type, to_obj.object_type_code)
        forward_relations.append({
            "relation_id": r.id,
            "relation_type": r.relation_type,
            "label": rt.get("name_zh", r.relation_type),
            "target_id": to_obj.id,
            "target_name": to_obj.name,
            "target_type": to_obj.object_type_code,
            "target_type_name": to_at.name if to_at else to_obj.object_type_code,
            "is_schema_valid": sk in schema_set,
        })

    # 反向关系
    rev_rels = db.query(OntologyRelation).filter(
        OntologyRelation.to_object_id == object_id,
        OntologyRelation.is_deleted == False,
    ).all()
    reverse_relations = []
    for r in rev_rels:
        from_obj = db.query(OntologyObject).filter(OntologyObject.id == r.from_object_id, OntologyObject.is_deleted == False).first()
        if not from_obj:
            continue
        from_at = db.query(AssetType).filter(AssetType.code == from_obj.object_type_code).first()
        rt = rel_type_map.get(r.relation_type, {})
        sk = (from_obj.object_type_code, r.relation_type, obj.object_type_code)
        reverse_relations.append({
            "relation_id": r.id,
            "relation_type": r.relation_type,
            "label": rt.get("reverse_name_zh", r.relation_type),
            "source_id": from_obj.id,
            "source_name": from_obj.name,
            "source_type": from_obj.object_type_code,
            "source_type_name": from_at.name if from_at else from_obj.object_type_code,
            "is_schema_valid": sk in schema_set,
        })

    return {
        "id": obj.id,
        "object_type_code": obj.object_type_code,
        "object_type_name": at.name if at else obj.object_type_code,
        "object_type_color": at.color if at else "#6366f1",
        "name": obj.name,
        "short_description": obj.short_description,
        "full_description": obj.full_description,
        "status": obj.status,
        "ai_confidence": float(obj.ai_confidence) if obj.ai_confidence else None,
        "tags": obj.tags_json or [],
        "attachment_filename": obj.attachment_filename,
        "attachment_path": obj.attachment_path,
        "external_link": obj.external_link,
        "kpi_formula": obj.kpi_formula,
        "created_by": obj.created_by,
        "created_at": obj.created_at.strftime("%Y-%m-%d %H:%M") if obj.created_at else None,
        "updated_at": obj.updated_at.strftime("%Y-%m-%d %H:%M") if obj.updated_at else None,
        "properties": properties,
        "forward_relations": forward_relations,
        "reverse_relations": reverse_relations,
    }


@app.get("/api/graph/stats")
def get_graph_stats(db: Session = Depends(get_db)):
    """获取图谱统计信息（用于筛选器计数展示）"""
    # 按类型统计
    type_stats = db.query(
        OntologyObject.object_type_code,
        func.count(OntologyObject.id)
    ).filter(OntologyObject.is_deleted == False).group_by(OntologyObject.object_type_code).all()

    # 按状态统计
    status_stats = db.query(
        OntologyObject.status,
        func.count(OntologyObject.id)
    ).filter(OntologyObject.is_deleted == False).group_by(OntologyObject.status).all()

    # 按关系类型统计
    rel_stats = db.query(
        OntologyRelation.relation_type,
        func.count(OntologyRelation.id)
    ).filter(OntologyRelation.is_deleted == False).group_by(OntologyRelation.relation_type).all()

    asset_type_map = _get_asset_type_map(db)
    rel_type_map = _get_relation_type_map(db)

    return {
        "type_stats": [{
            "code": code,
            "name": asset_type_map.get(code, {}).get("name", code),
            "color": asset_type_map.get(code, {}).get("color", "#6366f1"),
            "count": count,
        } for code, count in type_stats],
        "status_stats": [{"status": s, "count": c} for s, c in status_stats],
        "relation_stats": [{
            "code": code,
            "name": rel_type_map.get(code, {}).get("name_zh", code),
            "color": rel_type_map.get(code, {}).get("color", "#6b7280"),
            "count": count,
        } for code, count in rel_stats],
        "total_nodes": sum(c for _, c in type_stats),
        "total_edges": sum(c for _, c in rel_stats),
    }


# ============================================================
# 启动时补齐新增语义层表字段（兼容性处理）
# ============================================================

@app.on_event("startup")
def on_startup_semantic():
    """确保语义层新表存在（ORM 自动建表无法覆盖已有数据库的场景）"""
    # relation_types, relation_schema_rules, asset_type_property_schema 表已通过迁移脚本创建
    # ontology_relations 新字段已通过 ALTER TABLE 添加
    # 此处仅确保新 ORM 模型对应的表存在
    try:
        Base.metadata.create_all(bind=engine, checkfirst=True)
    except Exception:
        pass


# ============================================================
# 静态文件挂载（必须放在最后）
# ============================================================
app.mount("/static", StaticFiles(directory="static", html=True), name="static")
